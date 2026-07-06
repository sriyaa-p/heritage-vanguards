"""generate_pptx.py
Generates slide-01-mcp-server.pptx using python-pptx.
Run: pip install python-pptx && python generate_pptx.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Palette ──────────────────────────────────────────────────────────────────
BG_DARK    = RGBColor(0x02, 0x0d, 0x1f)
BG_CARD    = RGBColor(0x04, 0x1e, 0x2a)
CYAN       = RGBColor(0x00, 0xe5, 0xcc)
WHITE      = RGBColor(0xff, 0xff, 0xff)
GREY       = RGBColor(0x7a, 0xa8, 0xc0)
CODE_BG    = RGBColor(0x0b, 0x15, 0x20)
CODE_TEXT  = RGBColor(0xcd, 0xd9, 0xe5)
STEP_TEXT  = RGBColor(0x02, 0x0d, 0x1f)

# ── Slide size 16:9 1280x720 pt ───────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)  # 1280 / 96
prs.slide_height = Inches(7.5)     # 720  / 96

blank_layout = prs.slide_layouts[6]  # completely blank
slide = prs.slides.add_slide(blank_layout)
shapes = slide.shapes

def rgb_fill(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    rgb_fill(shape, color)
    return shape

def add_textbox(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb

# ── Background ───────────────────────────────────────────────────────────────
add_rect(slide, 0, 0, 13.333, 7.5, BG_DARK)

# ── Top badge (top-right): SLIDE 01 / 06 ─────────────────────────────────────
badge = add_rect(slide, 11.2, 0.35, 1.8, 0.28, BG_DARK)
badge.line.color.rgb = CYAN
badge.line.width = Pt(1)
add_textbox(slide, "SLIDE 01 / 06", 11.22, 0.36, 1.76, 0.26, 7.5, bold=True, color=CYAN, align=PP_ALIGN.CENTER)

# ── Title ─────────────────────────────────────────────────────────────────────
add_textbox(slide, "MCP Server ", 0.5, 0.38, 4.5, 0.6, 28, bold=True, color=WHITE)
add_textbox(slide, "Pipeline",    3.45, 0.38, 2.5, 0.6, 28, bold=True, color=CYAN)
add_textbox(slide, "HERITAGE VANGUARDS  —  AI-DRIVEN UNESCO NOMINATION ENGINE",
            0.5, 0.92, 10, 0.28, 8, color=GREY)

# ── Cards data ───────────────────────────────────────────────────────────────
cards = [
    (
        "STEP 01", "INGESTIONAGENT",
        "Standardizes varied submissions (text, coordinates, photos) into a canonical "
        "Pydantic dossier. Handles multilingual input via automatic language detection.",
        "LINGUA + GEMINI TRANSLATION"
    ),
    (
        "STEP 02", "REGISTRYAGENT",
        "Performs real-time database-level hybrid search. Gated: if FTS match score >= 0.3, "
        "Gemini checks semantics. Cuts duplicate submissions at the source.",
        "POSTGRESQL FTS + SEMANTIC SEARCH"
    ),
    (
        "STEP 03", "EVALUATIONAGENT",
        "Separates concerns. Evaluates text evidence across 8 UNESCO operational criteria. "
        "A deterministic keyword custom scorer ensures reproducible, auditable scores.",
        "DETERMINISTIC CUSTOM ENGINE"
    ),
    (
        "STEP 04", "MCP SERVER",
        "Exposes all 6 pipeline tools via Model Context Protocol. SSE transport on port 8001. "
        "Any MCP-compatible agent - Gemini, Claude, GPT - can orchestrate the full workflow.",
        "6 TOOLS  ·  STDIO + SSE TRANSPORT"
    ),
]

card_l = 0.5
card_w = 5.4
card_h = 1.25
card_top_start = 1.3
card_gap = 0.12

for i, (step, name, desc, tag) in enumerate(cards):
    top = card_top_start + i * (card_h + card_gap)
    # card background
    card = add_rect(slide, card_l, top, card_w, card_h, BG_CARD)
    card.line.color.rgb = RGBColor(0x00, 0x60, 0x58)
    card.line.width = Pt(0.75)
    # step badge
    badge_rect = add_rect(slide, card_l + 0.1, top + 0.1, 0.65, 0.2, CYAN)
    badge_rect.line.fill.background()
    add_textbox(slide, step, card_l + 0.1, top + 0.1, 0.65, 0.2, 6.5, bold=True, color=STEP_TEXT, align=PP_ALIGN.CENTER)
    # agent name
    add_textbox(slide, name, card_l + 0.1, top + 0.33, card_w - 0.2, 0.22, 10, bold=True, color=WHITE)
    # description
    add_textbox(slide, desc, card_l + 0.1, top + 0.55, card_w - 0.2, 0.52, 7.5, color=GREY)
    # tag
    tag_rect = add_rect(slide, card_l + 0.1, top + card_h - 0.27, len(tag)*0.065 + 0.2, 0.19, BG_DARK)
    tag_rect.line.color.rgb = RGBColor(0x00, 0x80, 0x70)
    tag_rect.line.width = Pt(0.75)
    add_textbox(slide, tag, card_l + 0.12, top + card_h - 0.27, len(tag)*0.065 + 0.2, 0.19, 6.5, bold=True, color=CYAN)

# ── Code panel ───────────────────────────────────────────────────────────────
code_l = 6.25
code_t = 1.3
code_w = 6.7
code_h = 5.55
code_box = add_rect(slide, code_l, code_t, code_w, code_h, CODE_BG)
code_box.line.color.rgb = RGBColor(0x00, 0x60, 0x58)
code_box.line.width = Pt(0.75)

# traffic lights
for xi, col in enumerate([RGBColor(0xff,0x5f,0x57), RGBColor(0xfe,0xbc,0x2e), RGBColor(0x28,0xc8,0x40)]):
    dot = slide.shapes.add_shape(9, Inches(code_l + 0.18 + xi*0.23), Inches(code_t + 0.14), Inches(0.13), Inches(0.13))
    rgb_fill(dot, col)
    dot.line.fill.background()

# filename label
add_textbox(slide, "mcp_server.py", code_l + 5.1, code_t + 0.1, 1.4, 0.22, 7.5, color=GREY)

# code lines
code_lines = [
    "# MCP Server - 6 Tools via SSE Transport",
    "from mcp.server.fastmcp import FastMCP",
    "import uvicorn",
    "",
    'app = FastMCP("heritage-pipeline")',
    "",
    "# Tool 1: Ingest & normalize submission",
    "@app.tool()",
    "async def ingest_submission(raw: dict) -> Dossier:",
    "    return await ingestion_agent.run(raw)",
    "",
    "# Tool 2: Registry deduplication check",
    "@app.tool()",
    "async def check_registry(dossier: Dossier):",
    "    score = await fts_search(dossier)",
    "    if score >= 0.3:",
    "        return await gemini.semantic_check(dossier)",
    "",
    "# Tool 3: Evaluate across 8 UNESCO criteria",
    "@app.tool()",
    "async def evaluate_criteria(dossier: Dossier):",
    "    return evaluation_agent.score(dossier)",
    "",
    "# SSE transport on port 8001",
    'if __name__ == "__main__":',
    '    uvicorn.run(app.sse_app(), host="0.0.0.0", port=8001)',
]

for li, line in enumerate(code_lines):
    lt = code_t + 0.42 + li * 0.19
    if lt > code_t + code_h - 0.1:
        break
    color = GREY if line.startswith("#") else CODE_TEXT
    add_textbox(slide, line, code_l + 0.18, lt, code_w - 0.3, 0.21, 7.5, color=color, italic=line.startswith("#"))

# ── Bottom stats bar ─────────────────────────────────────────────────────────
stats_t = 7.05
add_rect(slide, 0.5, stats_t - 0.02, 12.33, 0.01, RGBColor(0x00, 0x50, 0x48))

stats = [
    ("< 15 Minutes",  "Nominal Pipeline Runtime"),
    ("80% Savings",   "LLM Cost Cut via Deterministic Gating"),
    ("100% Accuracy", "Deterministic Scoring, 29 Tests Passing"),
]
stat_l = 0.5
for label, detail in stats:
    add_textbox(slide, label + "  ", stat_l, stats_t, 1.4, 0.35, 8, bold=True, color=CYAN)
    add_textbox(slide, detail, stat_l + 1.35, stats_t, 2.8, 0.35, 8, color=GREY)
    stat_l += 4.2

# ── Save ─────────────────────────────────────────────────────────────────────
out = "slide-01-mcp-server.pptx"
prs.save(out)
print(f"Saved: {out}")
